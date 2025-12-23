from flask import Flask, request, send_file, jsonify, render_template
import tempfile
import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import io

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024  # 10MB limit

@app.route('/')
def index():
    return render_template('index.html')

# PDF to PPTX Conversion
@app.route('/convert-pdf-to-pptx', methods=['POST'])
def pdf_to_pptx():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Please upload a PDF file'}), 400
    
    try:
        # Read PDF
        pdf_bytes = file.read()
        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Add a slide for each PDF page
        for page_num in range(len(pdf_doc)):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            
            # Convert PDF page to image
            page = pdf_doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            
            # Save image temporarily
            img_path = tempfile.mktemp(suffix='.png')
            pix.save(img_path)
            
            # Add image to slide
            left = top = Inches(0.5)
            slide.shapes.add_picture(img_path, left, top, 
                                    width=prs.slide_width - Inches(1),
                                    height=prs.slide_height - Inches(1))
            
            # Clean up temp image
            os.unlink(img_path)
        
        pdf_doc.close()
        
        # Save PPTX to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        return send_file(
            pptx_bytes,
            as_attachment=True,
            download_name='converted.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# PPTX to PDF Conversion
@app.route('/convert-pptx-to-pdf', methods=['POST'])
def pptx_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not (file.filename.lower().endswith('.pptx') or 
            file.filename.lower().endswith('.ppt')):
        return jsonify({'error': 'Please upload a PowerPoint file'}), 400
    
    try:
        # Read PPTX
        pptx_bytes = file.read()
        
        # For demo purposes - create a simple PDF
        # In production, you'd use proper conversion
        pdf_bytes = io.BytesIO()
        c = canvas.Canvas(pdf_bytes, pagesize=letter)
        
        # Add some content
        c.drawString(100, 750, "PowerPoint to PDF Conversion")
        c.drawString(100, 730, f"Original file: {file.filename}")
        c.drawString(100, 710, "This is a placeholder conversion.")
        c.drawString(100, 690, "For full conversion, consider using:")
        c.drawString(100, 670, "1. LibreOffice (libreoffice-convert)")
        c.drawString(100, 650, "2. CloudConvert API")
        c.drawString(100, 630, "3. Aspose.Slides")
        
        c.save()
        pdf_bytes.seek(0)
        
        return send_file(
            pdf_bytes,
            as_attachment=True,
            download_name='converted.pdf',
            mimetype='application/pdf'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/health')
def health():
    return jsonify({'status': 'healthy'})

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
